import sys
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from openpyxl import load_workbook
import aspose.slides as slides
import pptx
import colorsys

log = open('log.txt', mode='w', encoding="utf-8")
TEMPLATE = "master.pptx"
wb = load_workbook('./TeamsNerc2021.xlsx')
tx_left = tx_top = tx_width = tx_height = Inches(1)
sheet = wb['Sheet1']
prs = Presentation(TEMPLATE)
def get_layout(prs, layout_name):
    return next(x for x in prs.slide_layouts if x.name == layout_name)

def splice_team(info):
    i = 0
    while (info[i] != '('):
        i+=1
    team_name = info[:i]
    i+=1
    team_players = ""
    while (info[i]!= ')'):
        if (info[i] != ',' and info[i] != ' '):
            team_players+=info[i]
        elif (info[i] == ','):
            team_players += '\n'
        i+=1    
    return team_name, team_players
        
def add_slide(prs, layout_name, info_team, image_1, image_2, sheet, duration):
    layout = get_layout(prs, layout_name)
    slide = prs.slides.add_slide(layout)
    width = prs.slide_width
    height = prs.slide_height
    picture = slide.placeholders[10]
    picture.insert_picture(image_1)
    picture = slide.placeholders[12]
    picture.insert_picture(image_2)
    team_name, team_members = splice_team(info_team)
    title = slide.shapes.title
    title.text = team_name
    subtitle  = slide.placeholders[11]
    subtitle.text = team_members
    return slide

path = Path('./images')
list_of_files = []
for root, dirs, files in os.walk(path):
	for file in files:
		list_of_files.append(os.path.join(root,file))

for i in range (1, 116):
    img_path_1 = list_of_files[i * 2 - 1]
    img_path_2 = list_of_files[i * 2 - 2]
    cell = 'B' + str(i+1)
    layout_name = "Team Layout "
    if (i % 2 == 0):
        layout_name += "Left"
    else:
        layout_name += "Right"
    add_slide(prs, layout_name, sheet[cell].value, img_path_1, img_path_2, sheet, 2)
    print(i)
prs.save("teams_with_photo.pptx") 
log.close()




